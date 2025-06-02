import os
import pandas as pd
import fitz  
import docx
import io
import pptx
from django.shortcuts import render
from sentence_transformers import SentenceTransformer, util
from transformers import pipeline
from sklearn.metrics.pairwise import cosine_similarity

summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")
embedder = SentenceTransformer('all-MiniLM-L6-v2')


def extract_text(file):
    name = file.name.lower()

    if name.endswith('.txt'):
        return file.read().decode('utf-8')

    elif name.endswith('.docx'):
        doc = docx.Document(file)
        return '\n'.join([p.text for p in doc.paragraphs])

    elif name.endswith('.pdf'):
        text = ""
        with fitz.open(stream=file.read(), filetype="pdf") as pdf:
            for page in pdf:
                text += page.get_text()
        return text

    elif name.endswith('.pptx'):
        presentation = pptx.Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif name.endswith('.csv'):
        df = pd.read_csv(file)
        return df.to_string()

    elif name.endswith(('.xls', '.xlsx')):
        df = pd.read_excel(file)
        return df.to_string()

    else:
        return ""  


def summarize_sections(text):
    sections = {}
    current_section = "Introduction"
    buffer = []

    for line in text.splitlines():
        line = line.strip()
        if line.lower() in ["introduction", "methodology", "conclusion", "results", "discussion"]:
            if buffer:
                sections[current_section] = summarizer(' '.join(buffer[:20]))[0]['summary_text']
                buffer = []
            current_section = line
        else:
            buffer.append(line)
    if buffer:
        sections[current_section] = summarizer(' '.join(buffer[:20]))[0]['summary_text']

    return sections


def compare_docs(docs):
    summaries = [summarize_sections(doc) for doc in docs]
    all_keys = set(k for summary in summaries for k in summary.keys())
    comparison = {}

    for section in all_keys:
        section_texts = [summary.get(section, '') for summary in summaries]
        embeddings = embedder.encode(section_texts)
        if len(embeddings) >= 2:
            sim = cosine_similarity([embeddings[0]], embeddings[1:])[0]
            comparison[section] = {
                'texts': section_texts,
                'similarity': sim.tolist()
            }
        else:
            comparison[section] = {
                'texts': section_texts,
                'similarity': []
            }

    return comparison


def upload_view(request):
    if request.method == 'POST':
        files = request.FILES.getlist('documents')
        docs = [extract_text(f) for f in files]
        results = compare_docs(docs)
        return render(request, 'core/upload.html', {'results': results})
    return render(request, 'core/upload.html')